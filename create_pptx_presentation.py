#!/usr/bin/env python3
"""
Create a proper PowerPoint/LibreOffice Impress presentation
using python-pptx library for TaskFlow AWS Todo App
"""

import os
import sys

def create_pptx_presentation():
    """Create a PowerPoint presentation that can be opened in LibreOffice"""
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.text import MSO_ANCHOR
        from pptx.dml.color import RGBColor
    except ImportError:
        print("📦 Installing python-pptx...")
        os.system("pip install python-pptx")
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.text import MSO_ANCHOR
        from pptx.dml.color import RGBColor
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    def add_title_slide(prs):
        """Add title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "TaskFlow: AI-Powered Serverless Todo Application"
        subtitle.text = "Built on AWS with Infrastructure as Code\\n\\nAWS Cloud Computing Capstone Project\\nJuly 2025"
        
        # Style the title
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)  # Dark blue
        
        # Style the subtitle
        subtitle_para = subtitle.text_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(52, 73, 94)  # Darker blue
    
    def add_content_slide(prs, title_text, content_points, layout_type=1):
        """Add a content slide with bullet points"""
        slide_layout = prs.slide_layouts[layout_type]  # Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = title_text
        
        # Style title
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        
        # Add content
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(52, 73, 94)
            p.space_after = Pt(6)
    
    def add_two_column_slide(prs, title_text, left_content, right_content):
        """Add a two-column content slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title_text
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Left column
        left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        left_frame = left_shape.text_frame
        left_frame.clear()
        
        for i, point in enumerate(left_content):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(52, 73, 94)
            p.space_after = Pt(6)
        
        # Right column
        right_shape = slide.shapes.add_textbox(Inches(6.83), Inches(1.5), Inches(6), Inches(5.5))
        right_frame = right_shape.text_frame
        right_frame.clear()
        
        for i, point in enumerate(right_content):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(52, 73, 94)
            p.space_after = Pt(6)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "• Modern serverless todo application with AI capabilities",
        "• Complete AWS cloud architecture demonstration", 
        "• Cost-optimized within AWS Free Tier limits ($0/month)",
        "• Infrastructure as Code with Terraform",
        "• Production-ready with enterprise security",
        "",
        "Key Achievements:",
        "✅ Zero server management (100% serverless)",
        "✅ AI-powered task extraction from natural language", 
        "✅ Auto-scaling with pay-per-use pricing",
        "✅ Enterprise-grade authentication and security"
    ])
    
    # Slide 3: Problem Statement
    add_two_column_slide(prs, "Problem Statement", [
        "🔴 Traditional Todo Apps Limitations:",
        "",
        "• Server maintenance overhead",
        "• High infrastructure costs", 
        "• Poor scalability",
        "• Manual deployment processes",
        "• Security vulnerabilities",
        "• Limited intelligent features",
        "• Fixed capacity planning"
    ], [
        "🟢 TaskFlow Solution:",
        "",
        "• Serverless architecture eliminates servers",
        "• $0/month operating cost",
        "• Auto-scaling to millions of users",
        "• Automated CI/CD pipeline", 
        "• Enterprise-grade security",
        "• AI-powered task extraction",
        "• Pay-per-use pricing model"
    ])
    
    # Slide 4: AWS Architecture Overview
    add_content_slide(prs, "AWS Architecture Overview", [
        "Complete Serverless Stack:",
        "",
        "🎨 Frontend: Amplify, CloudFront, S3",
        "   • Static hosting, global CDN, asset storage",
        "",
        "🔐 Authentication: Cognito User Pools", 
        "   • User management, JWT tokens, MFA ready",
        "",
        "🌐 API Layer: API Gateway, Authorizers",
        "   • RESTful endpoints, security, CORS",
        "",
        "⚙️ Compute: Lambda Functions",
        "   • Python business logic, auto-scaling",
        "",
        "💾 Storage: DynamoDB",
        "   • NoSQL database, pay-per-request"
    ])
    
    # Slide 5: Cost Analysis
    add_content_slide(prs, "Cost Analysis - AWS Free Tier Optimization", [
        "💰 Monthly Operating Cost: $0.00",
        "",
        "Service Breakdown:",
        "• AWS Lambda: 1M requests/month (Current: ~10K) → $0.00",
        "• API Gateway: 1M API calls/month (Current: ~10K) → $0.00", 
        "• DynamoDB: 25 GB + 25 WCU/RCU (Current: ~1 GB) → $0.00",
        "• S3 (Amplify): 5 GB + 20K requests (Current: ~10 MB) → $0.00",
        "• Cognito: 50K monthly active users (Current: <100) → $0.00",
        "• CloudWatch: 10 metrics + 5 GB logs → $0.00",
        "",
        "💡 Free Tier Strategy Benefits:",
        "• Perfect for development and testing",
        "• Supports small to medium user bases",
        "• No upfront infrastructure investment"
    ])
    
    # Slide 6: Cost Scaling Projections
    add_content_slide(prs, "Cost Scaling Projections", [
        "Growth Scenarios - Pay-as-you-Scale:",
        "",
        "📊 Current (Free Tier): <100 users → $0.00/month",
        "",
        "📈 Small Business: 1,000 users → ~$9.00/month",
        "   • Lambda: $2.50, DynamoDB: $3.00, API Gateway: $3.50",
        "",
        "📈 Medium Business: 10,000 users → ~$90.00/month", 
        "   • Proportional scaling across all services",
        "",
        "📈 Enterprise: 100,000 users → ~$900.00/month",
        "   • Unlimited scaling potential",
        "",
        "💰 Cost Optimization Benefits:",
        "• No fixed infrastructure costs",
        "• Pay only for actual usage"
    ])
    
    # Slide 7: Innovation & Features
    add_content_slide(prs, "Innovation & Unique Features", [
        "🤖 AI-Powered Task Extraction:",
        "• Convert natural language to structured todos",
        "• Powered by Google Gemini API integration",
        "• Example: 'Meeting tomorrow at 3pm with client'",
        "",
        "🏗️ Serverless-First Architecture:",
        "• Zero cold start optimization (< 200ms)",
        "• Event-driven design patterns",
        "• Auto-scaling from 0 to 15,000+ concurrent users",
        "",
        "🔒 Security-First Approach:",
        "• JWT-based authentication with Amazon Cognito",
        "• API endpoint protection with authorizers",
        "• Least privilege IAM policies",
        "• Encryption at rest and in transit"
    ])
    
    # Slide 8: Technical Implementation
    add_two_column_slide(prs, "Technical Implementation", [
        "🎨 Frontend Technologies:",
        "• React 18 with Hooks",
        "• Tailwind CSS for styling", 
        "• AWS Amplify SDK",
        "• Responsive design principles",
        "",
        "⚙️ Backend Technologies:",
        "• Python 3.12 Lambda functions",
        "• Boto3 AWS SDK",
        "• RESTful API design",
        "• Comprehensive error handling"
    ], [
        "🏗️ Infrastructure as Code:",
        "• Terraform (AWS Provider v6.2.0)",
        "• Modular configuration",
        "• Environment management",
        "• Automated deployments",
        "",
        "🤖 AI Integration:",
        "• Google Gemini API",
        "• Natural language processing",
        "• Task extraction algorithms",
        "• Secure API key management"
    ])
    
    # Slide 9: Security Architecture
    add_content_slide(prs, "Security Architecture", [
        "🔒 Enterprise-Grade Security Implementation:",
        "",
        "Authentication & Authorization:",
        "• Amazon Cognito User Pools",
        "• JWT token validation",
        "• API Gateway authorizers",
        "",
        "Data Protection:",
        "• Encryption at rest (DynamoDB)",
        "• Encryption in transit (TLS/HTTPS)", 
        "• Secure API communications",
        "",
        "Access Control:",
        "• IAM roles and policies",
        "• Least privilege principles",
        "• Resource-based permissions"
    ])
    
    # Slide 10: Performance & Scalability
    add_content_slide(prs, "Performance & Scalability", [
        "🚀 Built for Scale from Day One:",
        "",
        "Auto-Scaling Capabilities:",
        "• Lambda: 0 → 15,000+ concurrent executions",
        "• DynamoDB: Auto-scaling WCU/RCU",
        "• API Gateway: Built-in traffic management",
        "• CloudFront: Global edge locations",
        "",
        "⚡ Performance Optimizations:",
        "• Lambda cold start minimization (< 200ms)",
        "• Efficient DynamoDB query patterns",
        "• CDN caching strategies",
        "• Optimized payload sizes",
        "• Connection pooling and reuse"
    ])
    
    # Slide 11: Key Features
    add_content_slide(prs, "Key Features Demo", [
        "🎯 Application Capabilities:",
        "",
        "👤 User Management:",
        "• User registration with email verification",
        "• Secure login/logout flow",
        "• Password reset functionality",
        "",
        "📝 Todo Operations:",
        "• Create, read, update, delete todos",
        "• Real-time synchronization",
        "• User-specific data isolation",
        "",
        "🤖 AI Features:",
        "• Natural language task extraction",
        "• Intelligent task parsing",
        "• Context-aware suggestions"
    ])
    
    # Slide 12: Development & DevOps
    add_two_column_slide(prs, "Development & DevOps", [
        "🛠️ Infrastructure as Code:",
        "• 100% Terraform managed",
        "• Version controlled infrastructure",
        "• Automated deployments",
        "• Environment consistency",
        "",
        "📋 Quality Assurance:",
        "• Comprehensive error handling",
        "• Input validation",
        "• Security best practices",
        "• Code documentation"
    ], [
        "🔄 Development Workflow:",
        "• Git-based version control",
        "• Automated testing pipeline",
        "• Continuous integration ready",
        "• Environment segregation",
        "",
        "🚀 Deployment Automation:",
        "• One-command deployment",
        "• Resource cleanup automation",
        "• Configuration validation"
    ])
    
    # Slide 13: Lessons Learned
    add_two_column_slide(prs, "Lessons Learned & Challenges", [
        "🔧 Technical Challenges:",
        "• Lambda cold start optimization",
        "• DynamoDB query pattern design",
        "• Cognito integration complexity",
        "• CORS configuration management",
        "",
        "💡 Solutions Implemented:",
        "• Lambda warming strategies",
        "• Efficient partition key design",
        "• Simplified authentication flow",
        "• Comprehensive CORS setup"
    ], [
        "🏆 Best Practices Discovered:",
        "• Infrastructure as Code benefits",
        "• Serverless monitoring importance",
        "• Security-first development",
        "• Cost optimization strategies",
        "",
        "📈 Key Learning Outcomes:",
        "• Cloud-native architecture patterns",
        "• AWS services integration",
        "• Modern development practices"
    ])
    
    # Slide 14: Future Enhancements
    add_content_slide(prs, "Future Enhancements", [
        "🔮 Roadmap & Scalability:",
        "",
        "📅 Short-term (3-6 months):",
        "• Real-time collaboration features with WebSockets",
        "• Mobile application development (React Native)",
        "• Advanced AI task categorization",
        "",
        "🚀 Medium-term (6-12 months):",
        "• Multi-region deployment for global users",
        "• Advanced analytics dashboard",
        "• Third-party integrations (Slack, Teams)",
        "",
        "🌟 Long-term (12+ months):",
        "• Machine learning recommendations",
        "• Enterprise SSO integration",
        "• API marketplace and plugin system"
    ])
    
    # Slide 15: Business Value & ROI
    add_content_slide(prs, "Business Value & ROI", [
        "💼 Project Impact & Benefits:",
        "",
        "💰 Financial ROI:",
        "• 100% elimination of server costs",
        "• 90% reduction in deployment time",
        "• Zero infrastructure maintenance",
        "",
        "🎯 Business Benefits:",
        "• $0 monthly operating costs",
        "• Production-ready architecture",
        "• Enterprise security standards",
        "",
        "🏆 Portfolio Value:",
        "• Demonstrates cloud expertise",
        "• Shows full-stack capabilities",
        "• Proves cost optimization skills"
    ])
    
    # Slide 16: Conclusion
    add_content_slide(prs, "Conclusion", [
        "🎯 Project Success Summary:",
        "",
        "✅ All Goals Successfully Achieved:",
        "• Complete serverless architecture implementation",
        "• AWS Free Tier cost optimization",
        "• Infrastructure as Code deployment",
        "• Production-ready application",
        "• AI-powered innovation features",
        "",
        "🚀 Technical Excellence:",
        "• Enterprise-grade security",
        "• Auto-scaling capabilities",
        "• Comprehensive monitoring",
        "• Modern development practices",
        "",
        "🎉 Ready for Production Deployment!"
    ])
    
    # Slide 17: Thank You
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add centered title
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Thank You!"
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 62, 80)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_shape.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "🎯 TaskFlow: AI-Powered Serverless Todo Application"
    subtitle_para.font.size = Pt(30)
    subtitle_para.font.color.rgb = RGBColor(52, 73, 94)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add contact info
    contact_shape = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.33), Inches(2))
    contact_frame = contact_shape.text_frame
    contact_para = contact_frame.paragraphs[0]
    contact_para.text = "Ready for Questions!\\n\\nGitHub Repository: AWS TaskFlow Todo Application\\nLive Demo: Production deployment available\\nDocumentation: Comprehensive setup guides included"
    contact_para.font.size = Pt(20)
    contact_para.font.color.rgb = RGBColor(52, 73, 94)
    contact_para.alignment = PP_ALIGN.CENTER
    
    return prs

def main():
    """Main function to create the presentation"""
    print("🚀 Creating TaskFlow AWS Todo App Presentation...")
    print("=" * 60)
    
    try:
        # Create the presentation
        prs = create_pptx_presentation()
        
        # Save the presentation
        output_file = "/home/zahurul/Documents/work/AWS_lab/capstone/aws_todo_app (Copy 2)/TaskFlow_AWS_Presentation.pptx"
        prs.save(output_file)
        
        print(f"✅ PowerPoint presentation created: {output_file}")
        print("📝 This file can be opened in LibreOffice Impress!")
        
        print("\n" + "=" * 60)
        print("✅ Presentation Creation Complete!")
        print("\n📋 Presentation includes:")
        print("   • Problem statement and solution overview")
        print("   • Complete AWS architecture details")
        print("   • Cost analysis and Free Tier optimization")
        print("   • Technical implementation details")
        print("   • Security and scalability features")
        print("   • Future roadmap and business value")
        print("\n💼 Ready for capstone project presentation!")
        
        return output_file
        
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        print("\n💡 Alternative: Use the HTML version created earlier")
        return None

if __name__ == "__main__":
    main()
